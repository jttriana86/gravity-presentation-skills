#!/usr/bin/env python3
"""
Gravity PPTX builder

Genera un .pptx en estilo Gravity desde un JSON spec.

Uso:
    python build_deck.py <input.json> <output.pptx>

Requisitos:
    pip install python-pptx
"""

import json
import sys
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

try:
    from PIL import Image as PILImage
except ImportError:
    PILImage = None


# ============================================================
# BRAND PALETTE
# ============================================================
PRIMARY    = RGBColor(0x05, 0x13, 0x67)
PRIMARY_2  = RGBColor(0x07, 0x16, 0x89)
BLACK      = RGBColor(0x00, 0x00, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x00, 0x47, 0x14)
GRAY       = RGBColor(0x66, 0x66, 0x66)
GRAY_LIGHT = RGBColor(0xE5, 0xE7, 0xEB)
RED_ALERT  = RGBColor(0xE7, 0x00, 0x39)
BG_SOFT    = RGBColor(0xF8, 0xF9, 0xFB)

# Soft white for navy backgrounds
WHITE_SOFT_60 = RGBColor(0xB3, 0xB3, 0xB3)
WHITE_SOFT_70 = RGBColor(0xC2, 0xC2, 0xC2)

# ============================================================
# FONTS
# ============================================================
FONT_HEADING = "Montserrat"
FONT_BODY    = "Open Sans"

# ============================================================
# SLIDE DIMENSIONS (16:9)
# ============================================================
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Padding internals
PAD_X = Inches(0.65)
PAD_Y = Inches(0.55)

# Vertical content bounds (entre header y page-num) para CENTRADO VERTICAL.
# El contenido de cada slide de cuerpo se agrupa al medio de estos límites.
CONTENT_TOP    = Inches(1.30)
CONTENT_BOTTOM = Inches(6.90)


def centered_start(content_h, top=CONTENT_TOP, bottom=CONTENT_BOTTOM):
    """Devuelve el Y (EMU) para que un bloque de altura content_h quede
    centrado verticalmente entre top y bottom. Nunca sube por encima de top."""
    avail = int(bottom) - int(top)
    return int(top) + max(0, (avail - int(content_h)) // 2)


# ============================================================
# HELPER UTILITIES
# ============================================================

def set_letter_spacing(run, points_x100):
    """Set letter-spacing on a run (units: 1/100 of a point)."""
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', str(points_x100))


def set_paragraph_line_spacing(paragraph, multiplier):
    """Set line spacing as a multiplier (0.95 = tight)."""
    paragraph.line_spacing = multiplier


def add_solid_bg(slide, color):
    """Add full-bleed solid background rectangle."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()  # no border
    bg.shadow.inherit = False
    return bg


def add_text_box(slide, left, top, width, height,
                 text, font_name=FONT_BODY, size=14,
                 bold=False, italic=False, color=BLACK,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 letter_spacing=None, line_spacing=None,
                 uppercase=False):
    """Add a text box with single-run text + formatting."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing

    text_to_use = text.upper() if uppercase else text
    run = p.add_run()
    run.text = text_to_use
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if letter_spacing is not None:
        set_letter_spacing(run, letter_spacing)

    return tb


def add_runs(paragraph, segments):
    """Add multiple runs with different styling to a single paragraph.
    segments: list of dicts {text, font, size, bold, italic, color, ...}
    """
    for seg in segments:
        run = paragraph.add_run()
        run.text = seg.get('text', '')
        run.font.name = seg.get('font', FONT_BODY)
        run.font.size = Pt(seg.get('size', 14))
        run.font.bold = seg.get('bold', False)
        run.font.italic = seg.get('italic', False)
        if seg.get('color'):
            run.font.color.rgb = seg['color']
        if seg.get('letter_spacing') is not None:
            set_letter_spacing(run, seg['letter_spacing'])


def parse_markdown_runs(text):
    """Parse a string with **bold** markers into segment list.
    Returns: list of (text, is_bold) tuples.
    """
    parts = re.split(r'(\*\*[^*]+\*\*)', text)
    segments = []
    for part in parts:
        if not part:
            continue
        if part.startswith('**') and part.endswith('**'):
            segments.append((part[2:-2], True))
        else:
            segments.append((part, False))
    return segments


def add_dot_seal(slide, left=PAD_X, top=PAD_Y):
    """Add the red dot brand seal (top-left corner)."""
    size = Pt(10)
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, size, size
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = RED_ALERT
    dot.line.fill.background()
    return dot


def add_brand_mark(slide, client_name, on_dark=False):
    """Add dot + brand name in top-left corner. Returns the textbox."""
    add_dot_seal(slide, PAD_X, PAD_Y + Pt(2))
    color = WHITE if on_dark else PRIMARY
    return add_text_box(
        slide,
        left=PAD_X + Pt(20),
        top=PAD_Y - Pt(2),
        width=Inches(4),
        height=Inches(0.3),
        text=client_name,
        font_name=FONT_HEADING,
        size=11,
        bold=True,
        color=color,
        letter_spacing=300,
        uppercase=True,
    )


def add_meta_top_right(slide, meta_text, on_dark=False):
    """Add small meta text in top-right corner."""
    color = WHITE_SOFT_60 if on_dark else GRAY
    return add_text_box(
        slide,
        left=SLIDE_W - PAD_X - Inches(5),
        top=PAD_Y,
        width=Inches(5),
        height=Inches(0.3),
        text=meta_text,
        font_name=FONT_BODY,
        size=10,
        bold=True,
        color=color,
        align=PP_ALIGN.RIGHT,
        letter_spacing=200,
        uppercase=True,
    )


def add_page_num(slide, current, total):
    """Add page number in bottom-right corner."""
    return add_text_box(
        slide,
        left=SLIDE_W - PAD_X - Inches(2),
        top=SLIDE_H - PAD_Y - Pt(8),
        width=Inches(2),
        height=Inches(0.3),
        text=f"{current:02d} / {total:02d}",
        font_name=FONT_HEADING,
        size=9,
        bold=True,
        color=GRAY,
        align=PP_ALIGN.RIGHT,
        letter_spacing=200,
    )


def add_card_bezel(slide, left, top, width, height):
    """Add a bezel-style card (white rounded rect + navy top accent + soft shadow).

    Returns: (card_shape, top_accent_shape) so caller can layer text on top.
    """
    # Main card body — rounded rectangle white
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.adjustments[0] = 0.04  # corner radius ~4%
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE0, 0xE3, 0xEA)
    card.line.width = Pt(0.5)

    # Soft shadow via XML manipulation
    spPr = card.fill._xPr
    sppr_xml = spPr
    # Add outer shadow effect
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    effectLst = etree.SubElement(sppr_xml, qn('a:effectLst'))
    outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '254000')  # 20pt blur
    outerShdw.set('dist', '50800')       # 4pt offset
    outerShdw.set('dir', '5400000')      # 90 degrees (down)
    outerShdw.set('algn', 'ctr')
    outerShdw.set('rotWithShape', '0')
    srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'))
    srgbClr.set('val', '051367')
    alpha = etree.SubElement(srgbClr, qn('a:alpha'))
    alpha.set('val', '15000')  # 15% opacity

    # Navy top accent bar (3pt thick)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = PRIMARY
    accent.line.fill.background()

    return card, accent


# ============================================================
# SLIDE BUILDERS
# ============================================================

def cover_slide(prs, data, current, total, client_name):
    """Build a cover slide (full navy + big white title)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    add_solid_bg(slide, PRIMARY)

    # Brand mark top-left (white version)
    add_brand_mark(slide, client_name, on_dark=True)

    # Center stack: eyebrow + title + subtitle
    cy = Inches(2.6)

    if data.get('eyebrow'):
        add_text_box(
            slide,
            left=PAD_X, top=cy,
            width=Inches(11), height=Inches(0.5),
            text=data['eyebrow'],
            font_name=FONT_HEADING, size=12, bold=True,
            color=WHITE_SOFT_70,
            letter_spacing=400,
            uppercase=True,
        )
        cy += Inches(0.55)

    add_text_box(
        slide,
        left=PAD_X, top=cy,
        width=Inches(11.5), height=Inches(2.5),
        text=data.get('title', ''),
        font_name=FONT_HEADING, size=72, bold=True,
        color=WHITE,
        letter_spacing=-200,
        line_spacing=0.95,
    )
    cy += Inches(2.2)

    if data.get('subtitle'):
        add_text_box(
            slide,
            left=PAD_X, top=cy,
            width=Inches(11), height=Inches(0.5),
            text=data['subtitle'],
            font_name=FONT_BODY, size=18, italic=True,
            color=WHITE_SOFT_70,
        )

    # Footer: logo (left) + meta (right)
    fy = SLIDE_H - PAD_Y - Inches(0.4)
    if data.get('footer_logo'):
        add_text_box(
            slide,
            left=PAD_X, top=fy,
            width=Inches(5), height=Inches(0.3),
            text=data['footer_logo'],
            font_name=FONT_HEADING, size=12, bold=True,
            color=WHITE,
            letter_spacing=200,
        )
    if data.get('footer_meta'):
        add_text_box(
            slide,
            left=SLIDE_W - PAD_X - Inches(6), top=fy,
            width=Inches(6), height=Inches(0.3),
            text=data['footer_meta'],
            font_name=FONT_BODY, size=10, bold=True,
            color=WHITE_SOFT_60,
            align=PP_ALIGN.RIGHT,
            letter_spacing=200,
            uppercase=True,
        )

    return slide


def divider_slide(prs, data, current, total, client_name):
    """Build a section divider slide (navy + huge outline number + title)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(slide, PRIMARY)
    add_brand_mark(slide, client_name, on_dark=True)

    # Big outline number (top of stack)
    cy = Inches(1.8)
    if data.get('number'):
        num_box = slide.shapes.add_textbox(
            PAD_X, cy, Inches(11), Inches(2.8)
        )
        tf = num_box.text_frame
        tf.margin_left = 0
        tf.margin_top = 0
        p = tf.paragraphs[0]
        p.line_spacing = 0.85
        run = p.add_run()
        run.text = data['number']
        run.font.name = FONT_HEADING
        run.font.size = Pt(180)
        run.font.bold = True
        # Outline effect (text fill = none, line = white 30%)
        rPr = run._r.get_or_add_rPr()
        # Remove default fill
        for child in rPr.findall(qn('a:solidFill')):
            rPr.remove(child)
        noFill = etree.SubElement(rPr, qn('a:noFill'))
        # Add line stroke
        ln = etree.SubElement(rPr, qn('a:ln'))
        ln.set('w', '19050')  # 1.5pt
        solidFill = etree.SubElement(ln, qn('a:solidFill'))
        srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgb.set('val', 'FFFFFF')
        alpha = etree.SubElement(srgb, qn('a:alpha'))
        alpha.set('val', '30000')

        cy += Inches(2.3)

    if data.get('title'):
        add_text_box(
            slide,
            left=PAD_X, top=cy,
            width=Inches(12), height=Inches(1.2),
            text=data['title'],
            font_name=FONT_HEADING, size=48, bold=True,
            color=WHITE,
            letter_spacing=-150,
            line_spacing=1.0,
        )
        cy += Inches(1.0)

    if data.get('subtitle'):
        add_text_box(
            slide,
            left=PAD_X, top=cy,
            width=Inches(11), height=Inches(0.5),
            text=data['subtitle'],
            font_name=FONT_BODY, size=16, italic=True,
            color=WHITE_SOFT_70,
        )

    return slide


def metrics_slide(prs, data, current, total, client_name):
    """Build a metrics slide (white bg + bezel cards + hero numbers + insight)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(slide, WHITE)

    # Header
    add_brand_mark(slide, client_name, on_dark=False)
    if data.get('meta'):
        add_meta_top_right(slide, data['meta'])

    # Title block — centrado vertical del grupo (título + cards)
    has_eb = bool(data.get('eyebrow'))
    has_sub = bool(data.get('subtitle'))
    title_h = (Inches(0.4) if has_eb else 0) + Inches(0.85) + (Inches(0.45) if has_sub else 0)
    # card_h debe cubrir todo lo que dibuja _populate_card (~3.08"): header, hero
    # metric, label, divider y la fila de stats. Con menos, los stats se salen.
    card_h = Inches(3.15)
    group_h = title_h + Inches(0.5) + card_h
    cy = centered_start(group_h, bottom=Inches(6.15))  # deja aire para el insight footer
    if data.get('eyebrow'):
        add_text_box(
            slide,
            left=PAD_X, top=cy,
            width=Inches(11), height=Inches(0.4),
            text=data['eyebrow'],
            font_name=FONT_HEADING, size=11, bold=True,
            color=GREEN,
            letter_spacing=400,
            uppercase=True,
        )
        cy += Inches(0.4)

    if data.get('title'):
        add_text_box(
            slide,
            left=PAD_X, top=cy,
            width=Inches(11), height=Inches(0.9),
            text=data['title'],
            font_name=FONT_HEADING, size=40, bold=True,
            color=PRIMARY,
            letter_spacing=-150,
            line_spacing=0.95,
        )
        cy += Inches(0.85)

    if data.get('subtitle'):
        add_text_box(
            slide,
            left=PAD_X, top=cy,
            width=Inches(11), height=Inches(0.4),
            text=data['subtitle'],
            font_name=FONT_BODY, size=14, italic=True,
            color=GRAY,
        )

    # Cards (card_top derivado del centrado del grupo: justo debajo del título)
    cards = data.get('cards', [])
    n = len(cards)
    card_top = cy + Inches(0.55)
    gap = Inches(0.25)
    available_w = SLIDE_W - 2 * PAD_X
    card_w = (available_w - gap * (n - 1)) / max(n, 1)

    for i, card_data in enumerate(cards):
        card_left = PAD_X + (card_w + gap) * i
        add_card_bezel(slide, card_left, card_top, card_w, card_h)
        _populate_card(slide, card_data, card_left, card_top, card_w, card_h)

    # Insight footer
    if data.get('insight'):
        ifoot_top = SLIDE_H - PAD_Y - Inches(0.7)
        # Tag
        tag_w = Inches(1.1)  # "INSIGHT" a 9pt con tracking no cabe en 0.9"
        tag = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, PAD_X, ifoot_top, tag_w, Inches(0.32)
        )
        tag.fill.solid()
        tag.fill.fore_color.rgb = PRIMARY
        tag.line.fill.background()
        tag_tf = tag.text_frame
        tag_tf.margin_left = Pt(8)
        tag_tf.margin_right = Pt(8)
        tag_tf.margin_top = Pt(2)
        tag_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tag_p = tag_tf.paragraphs[0]
        tag_p.alignment = PP_ALIGN.CENTER
        tag_run = tag_p.add_run()
        tag_run.text = "INSIGHT"
        tag_run.font.name = FONT_HEADING
        tag_run.font.size = Pt(9)
        tag_run.font.bold = True
        tag_run.font.color.rgb = WHITE
        set_letter_spacing(tag_run, 300)

        # Insight text with markdown
        insight_box = slide.shapes.add_textbox(
            PAD_X + tag_w + Pt(8), ifoot_top - Pt(2),
            SLIDE_W - 2 * PAD_X - tag_w - Pt(8), Inches(0.5)
        )
        itf = insight_box.text_frame
        itf.margin_left = 0
        itf.word_wrap = True
        itf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ip = itf.paragraphs[0]

        for text, is_bold in parse_markdown_runs(data['insight']):
            r = ip.add_run()
            r.text = text
            r.font.name = FONT_BODY
            r.font.size = Pt(13)
            if is_bold:
                r.font.bold = True
                r.font.italic = False
                r.font.color.rgb = PRIMARY
            else:
                r.font.italic = True
                r.font.color.rgb = BLACK

    add_page_num(slide, current, total)
    return slide


def _populate_card(slide, card_data, left, top, width, height):
    """Fill a bezel card with platform name, metric, stats."""
    inner_pad = Inches(0.3)
    cy = top + Pt(8) + inner_pad

    # Platform header row
    add_text_box(
        slide,
        left=left + inner_pad, top=cy,
        width=width - 2 * inner_pad - Inches(0.5), height=Inches(0.4),
        text=card_data.get('platform', ''),
        font_name=FONT_HEADING, size=12, bold=True,
        color=PRIMARY,
        letter_spacing=300,
        uppercase=True,
    )

    # Icon (small navy square with letter)
    icon_size = Inches(0.32)
    icon = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left + width - inner_pad - icon_size,
        cy + Pt(2),
        icon_size, icon_size
    )
    icon.adjustments[0] = 0.2
    icon.fill.solid()
    icon.fill.fore_color.rgb = PRIMARY
    icon.line.fill.background()
    if card_data.get('icon'):
        itf = icon.text_frame
        itf.margin_left = 0; itf.margin_right = 0
        itf.margin_top = 0; itf.margin_bottom = 0
        itf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ip = itf.paragraphs[0]
        ip.alignment = PP_ALIGN.CENTER
        ir = ip.add_run()
        ir.text = card_data['icon']
        ir.font.name = FONT_HEADING
        ir.font.size = Pt(10)
        ir.font.bold = True
        ir.font.color.rgb = WHITE

    cy += Inches(0.55)

    # Hero metric (arrow + percent)
    metric_text = card_data.get('metric', '')
    metric_box = slide.shapes.add_textbox(
        left + inner_pad, cy,
        width - 2 * inner_pad, Inches(0.9)
    )
    mtf = metric_box.text_frame
    mtf.margin_left = 0; mtf.word_wrap = False
    mp = mtf.paragraphs[0]
    arrow_run = mp.add_run()
    arrow_run.text = "▲ "
    arrow_run.font.name = FONT_HEADING
    arrow_run.font.size = Pt(20)
    arrow_run.font.bold = True
    arrow_run.font.color.rgb = GREEN
    metric_run = mp.add_run()
    metric_run.text = metric_text
    metric_run.font.name = FONT_HEADING
    metric_run.font.size = Pt(48)
    metric_run.font.bold = True
    metric_run.font.color.rgb = GREEN
    set_letter_spacing(metric_run, -150)

    cy += Inches(0.85)

    # Metric label
    if card_data.get('metric_label'):
        add_text_box(
            slide,
            left=left + inner_pad, top=cy,
            width=width - 2 * inner_pad, height=Inches(0.3),
            text=card_data['metric_label'],
            font_name=FONT_BODY, size=10, bold=True,
            color=GRAY,
            letter_spacing=200,
            uppercase=True,
        )
        cy += Inches(0.35)

    # Divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + inner_pad, cy,
        width - 2 * inner_pad, Pt(0.5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GRAY_LIGHT
    line.line.fill.background()
    cy += Inches(0.2)

    # Stats row
    stats = card_data.get('stats', [])
    if stats:
        stat_w = (width - 2 * inner_pad) / len(stats)
        for j, stat in enumerate(stats):
            sx = left + inner_pad + stat_w * j
            # Value
            add_text_box(
                slide,
                left=sx, top=cy,
                width=stat_w, height=Inches(0.4),
                text=str(stat.get('value', '')),
                font_name=FONT_HEADING, size=22, bold=True,
                color=BLACK,
            )
            # Label
            add_text_box(
                slide,
                left=sx, top=cy + Inches(0.42),
                width=stat_w, height=Inches(0.3),
                text=stat.get('label', ''),
                font_name=FONT_BODY, size=10,
                color=GRAY,
            )


def body_slide(prs, data, current, total, client_name):
    """Build a body slide (white bg + 2-column layout + numbered bullets)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(slide, WHITE)

    add_brand_mark(slide, client_name, on_dark=False)
    if data.get('meta'):
        add_meta_top_right(slide, data['meta'])

    # Centrado vertical: calcular alto del bloque de bullets para centrar
    bullet_h = Inches(0.95)
    _n = len(data.get('bullets', []))
    _intro_h = Inches(0.85) if data.get('intro') else 0
    content_h = _intro_h + _n * bullet_h
    start_y = centered_start(content_h)

    # Left column: eyebrow + title (40%)
    left_w = Inches(5.0)
    cy = start_y

    if data.get('eyebrow'):
        add_text_box(
            slide,
            left=PAD_X, top=cy,
            width=left_w, height=Inches(0.4),
            text=data['eyebrow'],
            font_name=FONT_HEADING, size=11, bold=True,
            color=GREEN,
            letter_spacing=400,
            uppercase=True,
        )
        cy += Inches(0.45)

    if data.get('title'):
        add_text_box(
            slide,
            left=PAD_X, top=cy,
            width=left_w, height=Inches(3),
            text=data['title'],
            font_name=FONT_HEADING, size=36, bold=True,
            color=PRIMARY,
            letter_spacing=-100,
            line_spacing=1.05,
        )

    # Right column: intro + bullets (60%)
    right_x = PAD_X + left_w + Inches(0.5)
    right_w = SLIDE_W - right_x - PAD_X
    ry = start_y

    if data.get('intro'):
        intro_box = slide.shapes.add_textbox(
            right_x, ry, right_w, Inches(0.8)
        )
        itf = intro_box.text_frame
        itf.margin_left = 0; itf.word_wrap = True
        ip = itf.paragraphs[0]
        ip.line_spacing = 1.4
        for text, is_bold in parse_markdown_runs(data['intro']):
            r = ip.add_run()
            r.text = text
            r.font.name = FONT_BODY
            r.font.size = Pt(14)
            if is_bold:
                r.font.bold = True
                r.font.color.rgb = PRIMARY
            else:
                r.font.color.rgb = BLACK
        ry += Inches(0.8)

    # Bullets
    bullets = data.get('bullets', [])
    bullet_h = Inches(0.95)
    for bullet in bullets:
        # Number circle
        num_size = Inches(0.4)
        num_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, right_x, ry, num_size, num_size
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = PRIMARY
        num_circle.line.fill.background()
        ntf = num_circle.text_frame
        ntf.margin_left = 0; ntf.margin_right = 0
        ntf.margin_top = 0; ntf.margin_bottom = 0
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        np_ = ntf.paragraphs[0]
        np_.alignment = PP_ALIGN.CENTER
        nr = np_.add_run()
        nr.text = bullet.get('number', '')
        nr.font.name = FONT_HEADING
        nr.font.size = Pt(11)
        nr.font.bold = True
        nr.font.color.rgb = WHITE

        # Title + description
        text_x = right_x + num_size + Inches(0.2)
        text_w = right_w - num_size - Inches(0.2)

        add_text_box(
            slide,
            left=text_x, top=ry - Pt(2),
            width=text_w, height=Inches(0.3),
            text=bullet.get('title', ''),
            font_name=FONT_BODY, size=14, bold=True,
            color=PRIMARY,
        )
        add_text_box(
            slide,
            left=text_x, top=ry + Inches(0.3),
            width=text_w, height=Inches(0.5),
            text=bullet.get('description', ''),
            font_name=FONT_BODY, size=12,
            color=GRAY,
            line_spacing=1.3,
        )
        ry += bullet_h

    add_page_num(slide, current, total)
    return slide


def image_slide(prs, data, current, total, client_name):
    """Build an image slide (white bg + compact title + large centered image).

    Use case: hand-drawn pizarra / sketchnote / diagram explainers.

    JSON keys:
      type: "image"
      eyebrow: short uppercase tag (optional)
      title: short title (kept compact at 32pt for max image space)
      meta: top-right meta text (optional)
      image: path to PNG/JPG (relative to JSON or absolute)
      caption: small italic caption below image (optional)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(slide, WHITE)

    add_brand_mark(slide, client_name, on_dark=False)
    if data.get('meta'):
        add_meta_top_right(slide, data['meta'])

    cy = Inches(1.05)
    if data.get('eyebrow'):
        add_text_box(
            slide,
            left=PAD_X, top=cy,
            width=Inches(11), height=Inches(0.35),
            text=data['eyebrow'],
            font_name=FONT_HEADING, size=11, bold=True,
            color=GREEN,
            letter_spacing=400,
            uppercase=True,
        )
        cy += Inches(0.4)

    if data.get('title'):
        add_text_box(
            slide,
            left=PAD_X, top=cy,
            width=Inches(11), height=Inches(0.85),
            text=data['title'],
            font_name=FONT_HEADING, size=32, bold=True,
            color=PRIMARY,
            letter_spacing=-100,
            line_spacing=1.0,
        )
        cy += Inches(0.85)

    # Image area: occupy remaining vertical space
    image_path = data.get('image')
    has_caption = bool(data.get('caption'))
    bottom_reserve = Inches(0.7) if has_caption else Inches(0.35)

    if image_path:
        img_path = Path(image_path)
        if not img_path.is_absolute():
            img_path = (Path(data.get('_spec_dir', '.')) / image_path).resolve()

        if not img_path.exists():
            print(f"  WARNING: image not found: {img_path}")
            return slide

        avail_x = PAD_X
        avail_y = cy + Inches(0.15)
        avail_w = SLIDE_W - 2 * PAD_X
        avail_h = SLIDE_H - avail_y - PAD_Y - bottom_reserve

        # Determine aspect-preserving fit using Pillow
        if PILImage is not None:
            with PILImage.open(str(img_path)) as im:
                iw, ih = im.size
            img_ratio = iw / ih
            avail_ratio = avail_w / avail_h
            if img_ratio > avail_ratio:
                final_w = avail_w
                final_h = int(avail_w / img_ratio)
            else:
                final_h = avail_h
                final_w = int(avail_h * img_ratio)
            final_x = avail_x + (avail_w - final_w) // 2
            final_y = avail_y + (avail_h - final_h) // 2
            slide.shapes.add_picture(str(img_path), final_x, final_y, final_w, final_h)
        else:
            # Fallback: let PowerPoint auto-size; place at top-left of available area
            slide.shapes.add_picture(str(img_path), avail_x, avail_y, avail_w, avail_h)

    if has_caption:
        add_text_box(
            slide,
            left=PAD_X,
            top=SLIDE_H - PAD_Y - Inches(0.55),
            width=SLIDE_W - 2 * PAD_X,
            height=Inches(0.35),
            text=data['caption'],
            font_name=FONT_BODY, size=12, italic=True,
            color=GRAY,
            align=PP_ALIGN.CENTER,
        )

    add_page_num(slide, current, total)
    return slide


def _add_soft_shadow(shape, color_hex='051367', alpha=10000, blur=190500, dist=38100):
    """Attach a subtle outer shadow to a shape via XML."""
    spPr = shape.fill._xPr
    effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', str(blur))
    outerShdw.set('dist', str(dist))
    outerShdw.set('dir', '5400000')
    outerShdw.set('algn', 'ctr')
    outerShdw.set('rotWithShape', '0')
    srgb = etree.SubElement(outerShdw, qn('a:srgbClr'))
    srgb.set('val', color_hex)
    a = etree.SubElement(srgb, qn('a:alpha'))
    a.set('val', str(alpha))


SITEMAP_SUBHEAD = RGBColor(0x8A, 0x96, 0xA3)   # gray-blue sub-section header
SITEMAP_FOOTER  = RGBColor(0x8A, 0x96, 0xA3)   # gray-blue footer pill
SITEMAP_BORDER  = RGBColor(0xD9, 0xDD, 0xE6)   # light pill border


def _sitemap_chip(slide, x, y, w, h, text, fill, txt_color, size, bold=True,
                  ls=150, radius=0.12, border=None, shadow_alpha=None):
    """Draw one rounded chip with centered text. Returns the shape."""
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  int(x), int(y), int(w), int(h))
    chip.adjustments[0] = radius
    chip.fill.solid()
    chip.fill.fore_color.rgb = fill
    if border is not None:
        chip.line.color.rgb = border
        chip.line.width = Pt(0.75)
    else:
        chip.line.fill.background()
    if shadow_alpha is not None:
        _add_soft_shadow(chip, alpha=shadow_alpha)
    else:
        chip.shadow.inherit = False
    tf = chip.text_frame
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 0.9
    r = p.add_run()
    r.text = text
    r.font.name = FONT_HEADING
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = txt_color
    if ls is not None:
        set_letter_spacing(r, ls)
    return chip


def sitemap_slide(prs, data, current, total, client_name):
    """Build a full sitemap slide (white bg + navy section headers + page pills).

    Gravity-styled site architecture map matching a reference IA diagram.
    Supports single columns and columns split into sub-columns, each with an
    optional gray "Footer" pill, and variable-height pills (text wraps).

    JSON keys:
      type: "sitemap"
      meta: top-right meta (optional)
      eyebrow: green uppercase tag (optional)
      title: navy title
      columns: list of either
        { "title": "Home", "pages": ["...", {"text":"...","highlight":true}],
          "footer": true }
        or (split section)
        { "title": "Institucional", "span": 2, "subcolumns": [
            { "title": "La empresa", "pages": [...], "footer": true }, ... ] }
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(slide, WHITE)
    add_brand_mark(slide, client_name, on_dark=False)
    if data.get('meta'):
        add_meta_top_right(slide, data['meta'])

    # Title block (compact, anchored top to maximize map space)
    cy = Inches(0.98)
    if data.get('eyebrow'):
        add_text_box(
            slide, left=PAD_X, top=cy, width=Inches(11), height=Inches(0.32),
            text=data['eyebrow'], font_name=FONT_HEADING, size=11, bold=True,
            color=GREEN, letter_spacing=400, uppercase=True,
        )
        cy += Inches(0.38)
    if data.get('title'):
        add_text_box(
            slide, left=PAD_X, top=cy, width=Inches(11.5), height=Inches(0.7),
            text=data['title'], font_name=FONT_HEADING, size=30, bold=True,
            color=PRIMARY, letter_spacing=-150, line_spacing=0.95,
        )
        cy += Inches(0.72)

    columns = data.get('columns', [])
    gap = Inches(0.13)
    avail_w = SLIDE_W - 2 * PAD_X
    slots = sum(c.get('span', 1) for c in columns) or 1
    slot_w = (int(avail_w) - int(gap) * (slots - 1)) / slots
    slot_pitch = slot_w + int(gap)

    header_y = int(cy + Inches(0.12))
    header_h = int(Inches(0.42))
    sub_y = header_y + header_h + int(Inches(0.13))
    subhead_h = int(Inches(0.32))
    pill_gap = int(Inches(0.075))
    pill_font = 8.5

    def slot_x(i):
        return int(PAD_X) + int(round(i * slot_pitch))

    def span_w(span):
        return int(round(slot_w * span + int(gap) * (span - 1)))

    def est_pill_h(text):
        usable_in = (slot_w / 914400.0) - 0.16
        char_w = 0.052 * (pill_font / 8.5)
        cpl = max(6, int(usable_in / char_w))
        lines = max(1, -(-len(text) // cpl))  # ceil division
        return int(Inches(0.26 + (lines - 1) * 0.135))

    def draw_pages(x, w, y, pages, footer_flag):
        for page in pages:
            if isinstance(page, dict):
                txt = page.get('text', ''); hl = page.get('highlight', False)
            else:
                txt = page; hl = False
            h = est_pill_h(txt)
            pill = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y), int(w), h)
            pill.adjustments[0] = 0.10
            pill.fill.solid()
            pill.shadow.inherit = False
            if hl:
                pill.fill.fore_color.rgb = PRIMARY
                pill.line.fill.background()
                txt_color = WHITE
            else:
                pill.fill.fore_color.rgb = WHITE
                pill.line.color.rgb = SITEMAP_BORDER
                pill.line.width = Pt(0.75)
                txt_color = PRIMARY
            ptf = pill.text_frame
            ptf.margin_left = Pt(4); ptf.margin_right = Pt(4)
            ptf.margin_top = Pt(1); ptf.margin_bottom = Pt(1)
            ptf.word_wrap = True
            ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
            pp = ptf.paragraphs[0]
            pp.alignment = PP_ALIGN.CENTER
            pp.line_spacing = 0.92
            pr = pp.add_run()
            pr.text = txt
            pr.font.name = FONT_BODY
            pr.font.size = Pt(pill_font)
            pr.font.color.rgb = txt_color
            y += h + pill_gap
        if footer_flag:
            fh = int(Inches(0.34))
            _sitemap_chip(slide, x, y, w, fh, "Footer", SITEMAP_FOOTER, WHITE,
                          size=9, bold=True, ls=200, radius=0.12)
            y += fh + pill_gap
        return y

    # Render columns left to right, advancing the slot cursor by span
    s = 0
    for col in columns:
        span = col.get('span', 1)
        x = slot_x(s)
        w = span_w(span)
        # Top-level navy header (spans all its slots)
        _sitemap_chip(slide, x, header_y, w, header_h, col.get('title', ''),
                      PRIMARY, WHITE, size=11, bold=True, ls=150,
                      shadow_alpha=18000)
        if col.get('subcolumns'):
            for k, sub in enumerate(col['subcolumns']):
                sx = slot_x(s + k)
                sw = int(round(slot_w))
                _sitemap_chip(slide, sx, sub_y, sw, subhead_h,
                              sub.get('title', ''), SITEMAP_SUBHEAD, WHITE,
                              size=9.5, bold=True, ls=100, radius=0.14)
                py = sub_y + subhead_h + int(Inches(0.10))
                draw_pages(sx, sw, py, sub.get('pages', []), sub.get('footer'))
        else:
            # Single column: pills start aligned with the sub-header row
            draw_pages(x, w, sub_y, col.get('pages', []), col.get('footer'))
        s += span

    add_page_num(slide, current, total)
    return slide


def quote_slide(prs, data, current, total, client_name):
    """Build a quote slide (white bg + giant italic quote + author)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(slide, WHITE)
    add_brand_mark(slide, client_name, on_dark=False)

    # Big quote mark
    add_text_box(
        slide,
        left=PAD_X, top=Inches(1.8),
        width=Inches(2), height=Inches(2),
        text='"',
        font_name=FONT_HEADING, size=180, bold=True,
        color=PRIMARY,
        line_spacing=0.8,
    )

    # Quote text (with optional highlight bold)
    text = data.get('text', '')
    highlight = data.get('highlight')

    quote_box = slide.shapes.add_textbox(
        PAD_X, Inches(2.6),
        SLIDE_W - 2 * PAD_X, Inches(3)
    )
    qtf = quote_box.text_frame
    qtf.margin_left = 0; qtf.word_wrap = True
    qp = qtf.paragraphs[0]
    qp.line_spacing = 1.15

    if highlight and highlight in text:
        before, after = text.split(highlight, 1)
        if before:
            r = qp.add_run()
            r.text = before
            r.font.name = FONT_BODY
            r.font.size = Pt(40)
            r.font.italic = True
            r.font.color.rgb = PRIMARY
        rh = qp.add_run()
        rh.text = highlight
        rh.font.name = FONT_BODY
        rh.font.size = Pt(40)
        rh.font.bold = True
        rh.font.italic = False
        rh.font.color.rgb = PRIMARY
        if after:
            r = qp.add_run()
            r.text = after
            r.font.name = FONT_BODY
            r.font.size = Pt(40)
            r.font.italic = True
            r.font.color.rgb = PRIMARY
    else:
        r = qp.add_run()
        r.text = text
        r.font.name = FONT_BODY
        r.font.size = Pt(40)
        r.font.italic = True
        r.font.color.rgb = PRIMARY

    # Author + role
    author_y = Inches(5.8)
    if data.get('author'):
        add_text_box(
            slide,
            left=PAD_X, top=author_y,
            width=Inches(8), height=Inches(0.4),
            text=f"— {data['author']}",
            font_name=FONT_HEADING, size=14, bold=True,
            color=BLACK,
            letter_spacing=100,
        )
    if data.get('role'):
        add_text_box(
            slide,
            left=PAD_X, top=author_y + Inches(0.35),
            width=Inches(8), height=Inches(0.3),
            text=data['role'],
            font_name=FONT_BODY, size=12, italic=True,
            color=GRAY,
        )

    add_page_num(slide, current, total)
    return slide


# ============================================================
# MAIN
# ============================================================

SLIDE_BUILDERS = {
    'cover':   cover_slide,
    'divider': divider_slide,
    'metrics': metrics_slide,
    'body':    body_slide,
    'quote':   quote_slide,
    'image':   image_slide,
    'sitemap': sitemap_slide,
}


def build(spec_path, output_path):
    """Read JSON spec, build .pptx, save."""
    spec_path = Path(spec_path)
    output_path = Path(output_path)

    with open(spec_path, 'r', encoding='utf-8') as f:
        spec = json.load(f)

    client_name = spec.get('client', 'CLIENTE')
    slides_spec = spec.get('slides', [])
    total = len(slides_spec)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print(f"Building deck: {client_name}")
    print(f"  Slides: {total}")

    spec_dir = str(spec_path.parent.resolve())

    for i, slide_data in enumerate(slides_spec, start=1):
        stype = slide_data.get('type', 'cover')
        builder = SLIDE_BUILDERS.get(stype)
        if not builder:
            print(f"  [{i}] WARNING: unknown slide type '{stype}', skipping")
            continue
        slide_data['_spec_dir'] = spec_dir
        builder(prs, slide_data, i, total, client_name)
        print(f"  [{i}] {stype}: {slide_data.get('title', slide_data.get('text', ''))[:50]}")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"\nDone -> {output_path.resolve()}")
    print(f"Size: {output_path.stat().st_size / 1024:.1f} KB")


if __name__ == '__main__':
    if len(sys.argv) != 3:
        print("Usage: python build_deck.py <input.json> <output.pptx>")
        sys.exit(1)
    build(sys.argv[1], sys.argv[2])
